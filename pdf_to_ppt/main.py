"""
PDF -> Premium PowerPoint  |  Google Gemini
Usage:  python main.py <input.pdf> [output.pptx]
"""

import os, sys, json, re, time
import fitz
from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ══════════════════════════════════════════════════════════════════════════════
# PALETTE  — deep space + electric violet + neon mint
# ══════════════════════════════════════════════════════════════════════════════
BG          = RGBColor(0x08, 0x0C, 0x14)   # near-black
SURFACE     = RGBColor(0x10, 0x18, 0x28)   # card base
SURFACE2    = RGBColor(0x18, 0x24, 0x38)   # lighter card
VIOLET      = RGBColor(0x7C, 0x3A, 0xED)   # electric violet
VIOLET_SOFT = RGBColor(0xA7, 0x8B, 0xFA)   # soft violet
MINT        = RGBColor(0x10, 0xB9, 0x81)   # neon mint/teal
MINT_SOFT   = RGBColor(0x6E, 0xE7, 0xB7)   # pale mint
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)   # warm amber
ROSE        = RGBColor(0xF4, 0x3F, 0x5E)   # rose accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE    = RGBColor(0xE2, 0xE8, 0xF0)
MUTED       = RGBColor(0x64, 0x74, 0x8B)
DIVIDER     = RGBColor(0x1E, 0x2D, 0x45)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Accent colour per slide index (cycles)
ACCENTS = [VIOLET, MINT, AMBER, ROSE, VIOLET_SOFT, MINT_SOFT]

# Icon glyphs
ICONS = {
    "code":   "{ }",  "bug":    "⚡",  "tools":  "◎",
    "chart":  "▲",    "book":   "≡",   "check":  "✦",
    "rocket": "⟶",   "shield": "◈",
}

TEXT_MODELS = [
    "models/gemini-2.0-flash-lite",
    "models/gemini-2.0-flash-001",
    "models/gemini-2.0-flash",
    "models/gemini-2.5-flash",
    "models/gemini-flash-latest",
]


# ══════════════════════════════════════════════════════════════════════════════
# PDF EXTRACTION
# ══════════════════════════════════════════════════════════════════════════════
def extract_pdf_text(pdf_path: str) -> str:
    doc = fitz.open(pdf_path)
    if doc.page_count == 0:
        raise ValueError("PDF has no pages.")
    pages = [p.get_text("text").strip() for p in doc if p.get_text("text").strip()]
    doc.close()
    text = "\n\n".join(pages)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    return text.strip()


# ══════════════════════════════════════════════════════════════════════════════
# GEMINI ANALYSIS
# ══════════════════════════════════════════════════════════════════════════════
PROMPT = """You are a senior presentation designer. Analyze the document and return ONLY valid JSON.

Required structure:
{{
  "title": "concise document title",
  "subtitle": "one compelling line describing the topic",
  "summary": "2-3 sentence executive summary",
  "tag": "2-word category tag e.g. SOFTWARE ENGINEERING",
  "slides": [
    {{
      "heading": "slide title (max 5 words)",
      "icon": "one of: code|bug|tools|chart|book|check|rocket|shield",
      "points": ["concise point, max 12 words", "3 to 5 points"],
      "stat": "one striking short fact or metric, e.g. '3x faster' or '40% less bugs'",
      "stat_label": "2-3 word label for the stat, e.g. 'Performance Gain'"
    }}
  ]
}}

Rules:
- 4 to 6 slides
- Headings max 5 words, punchy
- Bullets max 12 words, action-oriented
- stat is empty string if nothing meaningful available
- Pure JSON only, no markdown, no code fences

Document:
{text}
"""


def call_gemini(text: str) -> dict:
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        raise EnvironmentError(
            "GEMINI_API_KEY not set.\n"
            "  Windows:  set GEMINI_API_KEY=your_key\n"
            "  Mac/Linux: export GEMINI_API_KEY=your_key"
        )
    client = genai.Client(api_key=api_key)
    prompt = PROMPT.format(text=text[:12000])
    raw = None

    for model in TEXT_MODELS:
        try:
            print(f"      [{model}] ...", end=" ", flush=True)
            r = client.models.generate_content(model=model, contents=prompt)
            raw = r.text.strip()
            print("OK")
            break
        except Exception as e:
            msg = str(e)
            if any(x in msg for x in ["429", "RESOURCE_EXHAUSTED", "quota"]):
                print("quota exhausted")
            elif any(x in msg for x in ["503", "UNAVAILABLE"]):
                print("overloaded, waiting 15s...")
                time.sleep(15)
                try:
                    r = client.models.generate_content(model=model, contents=prompt)
                    raw = r.text.strip(); print("OK (retry)"); break
                except Exception:
                    pass
            elif any(x in msg for x in ["404", "NOT_FOUND"]):
                print("not found")
            else:
                raise

    if raw is None:
        raise RuntimeError(
            "All Gemini models quota-exhausted.\n"
            "Get a new key: https://aistudio.google.com/apikey"
        )

    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```\s*$", "", raw)
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        raise ValueError(f"Gemini returned invalid JSON:\n{raw}\n{e}")


# ══════════════════════════════════════════════════════════════════════════════
# DRAWING PRIMITIVES
# ══════════════════════════════════════════════════════════════════════════════
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, l, t, w, h, color: RGBColor, alpha: int = None):
    """Draw a filled rectangle. alpha 0-100000 (pptx units, 100000=opaque)."""
    sp = slide.shapes.add_shape(1, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    if alpha is not None:
        # inject transparency via XML
        xfrm = sp._element
        solidFill = xfrm.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                a_el = etree.SubElement(srgb, qn('a:alpha'))
                a_el.set('val', str(alpha))
    return sp

def _circle(slide, l, t, d, color: RGBColor):
    sp = slide.shapes.add_shape(9, l, t, d, d)   # 9 = oval
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    return sp

def _tb(slide, txt, l, t, w, h, size=16, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = txt
    run.font.name      = font
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    return box

def _multiline_tb(slide, lines: list, l, t, w, h,
                  size=14, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, line_space_pt=6):
    """Textbox with multiple paragraphs."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(line_space_pt)
        run = p.add_run()
        run.text           = line
        run.font.name      = "Calibri"
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.color.rgb = color
    return box


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE  (full-bleed split design)
# ══════════════════════════════════════════════════════════════════════════════
def build_title_slide(prs, title, subtitle, summary, tag):
    slide = _blank(prs)
    _bg(slide, BG)

    # ── background geometry ──────────────────────────────────────────────────

    # Right panel (60% width) — slightly lighter
    _rect(slide, Inches(5.3), 0, Inches(8.03), SLIDE_H, SURFACE)

    # Diagonal slash divider (simulate with rotated rect via overlapping rects)
    # We use a parallelogram-like effect with two triangles
    # Top-right large circle cluster
    _circle(slide, Inches(9.8),  Inches(-1.5), Inches(5.0), SURFACE2)
    _circle(slide, Inches(11.0), Inches(-0.8), Inches(3.2), DIVIDER)
    _circle(slide, Inches(10.5), Inches(0.3),  Inches(1.8), VIOLET)

    # Bottom-right small circles
    _circle(slide, Inches(12.2), Inches(6.2),  Inches(1.8), SURFACE2)
    _circle(slide, Inches(12.8), Inches(6.6),  Inches(1.0), MINT)

    # Bottom-left accent bar
    _rect(slide, 0, Inches(7.1), Inches(5.3), Inches(0.4), VIOLET)

    # Thin top border line
    _rect(slide, 0, 0, SLIDE_W, Inches(0.04), VIOLET)

    # Left panel vertical accent stripe
    _rect(slide, Inches(0.38), Inches(0.6), Inches(0.05), Inches(5.6), VIOLET)

    # ── left panel text ──────────────────────────────────────────────────────

    # Tag pill
    _rect(slide, Inches(0.6), Inches(0.65), Inches(2.4), Inches(0.35), VIOLET)
    _tb(slide, tag or "RESEARCH PAPER",
        Inches(0.62), Inches(0.67), Inches(2.36), Inches(0.32),
        size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Main title — large, bold
    _tb(slide, title,
        Inches(0.6), Inches(1.2), Inches(4.5), Inches(3.0),
        size=34, bold=True, color=WHITE)

    # Mint divider line
    _rect(slide, Inches(0.6), Inches(4.35), Inches(1.8), Inches(0.04), MINT)

    # Subtitle
    _tb(slide, subtitle,
        Inches(0.6), Inches(4.5), Inches(4.5), Inches(0.65),
        size=13, color=MINT_SOFT, italic=True)

    # Branding
    _tb(slide, "Generated with Google Gemini AI",
        Inches(0.6), Inches(7.15), Inches(4.0), Inches(0.3),
        size=8, color=MUTED)

    # ── right panel text ─────────────────────────────────────────────────────

    # "OVERVIEW" label
    _tb(slide, "OVERVIEW",
        Inches(5.7), Inches(1.0), Inches(7.0), Inches(0.4),
        size=9, bold=True, color=VIOLET_SOFT)

    # Summary heading
    _tb(slide, "Executive Summary",
        Inches(5.7), Inches(1.5), Inches(7.2), Inches(0.7),
        size=22, bold=True, color=WHITE)

    # Mint underline
    _rect(slide, Inches(5.7), Inches(2.28), Inches(1.2), Inches(0.04), MINT)

    # Summary body
    _tb(slide, summary,
        Inches(5.7), Inches(2.45), Inches(7.2), Inches(3.5),
        size=13, color=OFFWHITE)

    # Bottom right decoration
    _tb(slide, "↓  Scroll to explore",
        Inches(5.7), Inches(6.5), Inches(4.0), Inches(0.4),
        size=9, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2–N — CONTENT  (3-zone layout)
# ══════════════════════════════════════════════════════════════════════════════
def build_content_slide(prs, heading, points, icon, stat, stat_label,
                        slide_num, total, accent: RGBColor):
    slide = _blank(prs)
    _bg(slide, BG)

    # ── background layers ────────────────────────────────────────────────────

    # Thin top accent bar
    _rect(slide, 0, 0, SLIDE_W, Inches(0.05), accent)

    # Header band
    _rect(slide, 0, 0, SLIDE_W, Inches(1.1), SURFACE)

    # Right sidebar
    _rect(slide, Inches(9.6), Inches(1.1), Inches(3.73), Inches(6.4), SURFACE)

    # Decorative large circle (top-right, partially clipped)
    _circle(slide, Inches(10.5), Inches(-0.8), Inches(3.8), SURFACE2)
    _circle(slide, Inches(11.4), Inches(-0.2), Inches(2.2), DIVIDER)

    # Small accent dot cluster (bottom-left)
    _circle(slide, Inches(0.1),  Inches(6.8), Inches(0.5), SURFACE2)
    _circle(slide, Inches(0.45), Inches(6.9), Inches(0.3), accent)

    # ── header ───────────────────────────────────────────────────────────────

    # Icon badge (circle)
    _circle(slide, Inches(0.2), Inches(0.15), Inches(0.75), accent)
    icon_char = ICONS.get(icon.lower(), "◆")
    _tb(slide, icon_char,
        Inches(0.2), Inches(0.22), Inches(0.75), Inches(0.6),
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Heading
    _tb(slide, heading,
        Inches(1.1), Inches(0.15), Inches(8.2), Inches(0.8),
        size=28, bold=True, color=WHITE)

    # Slide counter pill
    _rect(slide, Inches(12.0), Inches(0.25), Inches(1.0), Inches(0.38), accent)
    _tb(slide, f"{slide_num}  /  {total}",
        Inches(12.0), Inches(0.27), Inches(1.0), Inches(0.34),
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Accent underline below header
    _rect(slide, 0, Inches(1.1), Inches(9.6), Inches(0.03), accent)

    # ── bullet cards (left+center zone) ──────────────────────────────────────
    card_l  = Inches(0.25)
    card_w  = Inches(9.1)
    card_h  = Inches(0.98)
    gap     = Inches(0.13)
    start_t = Inches(1.22)

    for i, point in enumerate(points[:5]):
        t   = start_t + i * (card_h + gap)
        bg  = SURFACE2 if i % 2 == 0 else SURFACE

        # Card
        _rect(slide, card_l, t, card_w, card_h, bg)

        # Left accent stripe
        _rect(slide, card_l, t, Inches(0.04), card_h, accent)

        # Number circle
        _circle(slide, card_l + Inches(0.1), t + Inches(0.24), Inches(0.5), accent)
        _tb(slide, str(i + 1),
            card_l + Inches(0.1), t + Inches(0.26),
            Inches(0.5), Inches(0.45),
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Point text
        _tb(slide, point,
            card_l + Inches(0.72), t + Inches(0.16),
            card_w - Inches(0.82), card_h - Inches(0.32),
            size=14, color=OFFWHITE)

    # ── right sidebar ─────────────────────────────────────────────────────────

    # Large icon display
    _circle(slide, Inches(10.15), Inches(1.4), Inches(2.4), SURFACE2)
    _tb(slide, icon_char,
        Inches(10.15), Inches(2.1), Inches(2.4), Inches(1.1),
        size=44, bold=True, color=accent, align=PP_ALIGN.CENTER)

    # Stat card
    if stat and stat.strip():
        _rect(slide, Inches(9.75), Inches(4.1), Inches(3.3), Inches(1.8), accent)

        # Decorative corner dot
        _circle(slide, Inches(12.7), Inches(4.0), Inches(0.5), SURFACE2)

        _tb(slide, stat,
            Inches(9.8), Inches(4.2), Inches(3.2), Inches(0.9),
            size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tb(slide, (stat_label or "KEY INSIGHT").upper(),
            Inches(9.8), Inches(5.05), Inches(3.2), Inches(0.4),
            size=9, bold=True, color=BG, align=PP_ALIGN.CENTER)
    else:
        # Decorative dots grid when no stat
        for row in range(3):
            for col in range(4):
                _circle(slide,
                        Inches(9.9 + col * 0.55),
                        Inches(4.2 + row * 0.55),
                        Inches(0.18),
                        SURFACE2 if (row + col) % 2 == 0 else DIVIDER)

    # Slide label (bottom of sidebar)
    _tb(slide, f"SLIDE  {slide_num:02d}",
        Inches(9.75), Inches(6.3), Inches(3.3), Inches(0.4),
        size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    # ── footer ────────────────────────────────────────────────────────────────
    _rect(slide, 0, Inches(7.18), SLIDE_W, Inches(0.32), SURFACE)
    _rect(slide, 0, Inches(7.18), Inches(0.04), Inches(0.32), accent)
    _tb(slide, heading.upper(),
        Inches(0.2), Inches(7.2), Inches(9.0), Inches(0.28),
        size=8, color=MUTED)
    _tb(slide, "Google Gemini AI",
        Inches(11.5), Inches(7.2), Inches(1.7), Inches(0.28),
        size=8, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# LAST SLIDE — SUMMARY  (bento-grid layout)
# ══════════════════════════════════════════════════════════════════════════════
def build_summary_slide(prs, slides_data, title):
    slide = _blank(prs)
    _bg(slide, BG)

    # Top accent bar
    _rect(slide, 0, 0, SLIDE_W, Inches(0.05), VIOLET)

    # Header band
    _rect(slide, 0, 0, SLIDE_W, Inches(1.35), SURFACE)

    # Decorative circles
    _circle(slide, Inches(11.2), Inches(-0.6), Inches(3.0), SURFACE2)
    _circle(slide, Inches(12.1), Inches(-0.1), Inches(1.6), DIVIDER)
    _circle(slide, Inches(12.5), Inches(0.5),  Inches(0.8), VIOLET)

    # Header text
    _tb(slide, "AGENDA  /  SUMMARY",
        Inches(0.4), Inches(0.1), Inches(6.0), Inches(0.4),
        size=9, bold=True, color=VIOLET_SOFT)
    _tb(slide, "What We Covered",
        Inches(0.4), Inches(0.5), Inches(8.0), Inches(0.75),
        size=28, bold=True, color=WHITE)

    # Accent underline
    _rect(slide, Inches(0.4), Inches(1.32), Inches(1.5), Inches(0.04), MINT)

    # ── bento grid ────────────────────────────────────────────────────────────
    items   = slides_data[:6]
    n       = len(items)
    cols    = 3
    rows    = (n + cols - 1) // cols
    cell_w  = Inches(4.1)
    cell_h  = Inches(1.7) if rows <= 2 else Inches(1.4)
    gap_x   = Inches(0.25)
    gap_y   = Inches(0.2)
    start_l = Inches(0.35)
    start_t = Inches(1.55)

    for idx, s in enumerate(items):
        col = idx % cols
        row = idx // cols
        l   = start_l + col * (cell_w + gap_x)
        t   = start_t + row * (cell_h + gap_y)
        acc = ACCENTS[idx % len(ACCENTS)]

        # Cell background
        _rect(slide, l, t, cell_w, cell_h, SURFACE)

        # Top accent bar on cell
        _rect(slide, l, t, cell_w, Inches(0.04), acc)

        # Number
        _tb(slide, f"{idx+1:02d}",
            l + Inches(0.12), t + Inches(0.08),
            Inches(0.5), Inches(0.38),
            size=11, bold=True, color=acc)

        # Icon
        icon_char = ICONS.get(s.get("icon", "code").lower(), "◆")
        _tb(slide, icon_char,
            l + cell_w - Inches(0.55), t + Inches(0.08),
            Inches(0.45), Inches(0.38),
            size=14, bold=True, color=acc, align=PP_ALIGN.RIGHT)

        # Heading
        _tb(slide, s["heading"],
            l + Inches(0.12), t + Inches(0.5),
            cell_w - Inches(0.24), cell_h - Inches(0.65),
            size=13, bold=True, color=OFFWHITE)

        # Bottom accent line
        _rect(slide, l, t + cell_h - Inches(0.04), cell_w, Inches(0.04), acc)

    # ── footer ────────────────────────────────────────────────────────────────
    _rect(slide, 0, Inches(7.18), SLIDE_W, Inches(0.32), SURFACE)
    _rect(slide, 0, Inches(7.18), Inches(0.04), Inches(0.32), VIOLET)
    _tb(slide, title,
        Inches(0.2), Inches(7.2), Inches(10.0), Inches(0.28),
        size=8, color=MUTED)
    _tb(slide, "Google Gemini AI",
        Inches(11.5), Inches(7.2), Inches(1.7), Inches(0.28),
        size=8, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# ASSEMBLE
# ══════════════════════════════════════════════════════════════════════════════
def build_pptx(data: dict, output_path: str):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = data.get("slides", [])
    total  = len(slides)

    build_title_slide(
        prs,
        title    = data["title"],
        subtitle = data.get("subtitle", ""),
        summary  = data["summary"],
        tag      = data.get("tag", "RESEARCH PAPER"),
    )

    for i, s in enumerate(slides):
        build_content_slide(
            prs,
            heading    = s["heading"],
            points     = s["points"],
            icon       = s.get("icon", "code"),
            stat       = s.get("stat", ""),
            stat_label = s.get("stat_label", ""),
            slide_num  = i + 1,
            total      = total,
            accent     = ACCENTS[i % len(ACCENTS)],
        )

    build_summary_slide(prs, slides, data["title"])
    prs.save(output_path)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main():
    if len(sys.argv) < 2:
        print("Usage: python main.py <input.pdf> [output.pptx]")
        sys.exit(1)

    pdf_path    = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else "final_presentation.pptx"

    if not os.path.isfile(pdf_path):
        print(f"Error: '{pdf_path}' not found.")
        sys.exit(1)

    print("\n[1/3] Extracting PDF text...")
    text = extract_pdf_text(pdf_path)
    print(f"      {len(text):,} characters extracted.")

    print("\n[2/3] Gemini AI analysis...")
    data = call_gemini(text)
    print(f"      Title   : {data.get('title')}")
    print(f"      Subtitle: {data.get('subtitle')}")
    print(f"      Slides  : {len(data.get('slides', []))}")

    print(f"\n[3/3] Building PowerPoint -> '{output_path}'...")
    build_pptx(data, output_path)
    size_kb = os.path.getsize(output_path) // 1024
    print(f"\n✅  Done!  {output_path}  ({size_kb} KB)")
    print(f"    {len(data['slides'])+2} slides total  "
          f"(1 title + {len(data['slides'])} content + 1 summary)")


if __name__ == "__main__":
    main()
